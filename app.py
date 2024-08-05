import sqlite3
import PyPDF2
import pytesseract
from pdf2image import convert_from_path
import streamlit as st
from pptx import Presentation
from docx import Document
from PIL import Image
from langchain_community.embeddings import OllamaEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain.chains import ConversationalRetrievalChain
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain.memory import ConversationBufferMemory
from langchain_groq import ChatGroq


# Set up page configuration's
st.set_page_config(page_icon="📄", layout="wide", page_title="PDF Conversational AI")

# Access the API key from Streamlit secrets
try:
    groq_api_key = st.secrets["groq"]["api_key"]
except KeyError:
    st.error("API key for GROQ is missing. Please check your secrets configuration.")
    st.stop()

# Initialize GROQ chat with the provided API key, model name, and settings
try:
    llm_groq = ChatGroq(
        groq_api_key=groq_api_key,
        model_name="Gemma2-9b-it",
        temperature=0.2
    )
except Exception as e:
    st.error(f"Failed to initialize GROQ chat: {e}")
    st.stop()

# Initialize the SQLite database
conn = sqlite3.connect('file_history.db')
c = conn.cursor()

# Create or alter tables if they do not exist
c.execute('''CREATE TABLE IF NOT EXISTS uploaded_files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_name TEXT UNIQUE
            )''')

c.execute('''CREATE TABLE IF NOT EXISTS chat_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                question TEXT,
                answer TEXT
            )''')
conn.commit()

def process_pdf(file):
    pdf_text = ""
    try:
        pdf = PyPDF2.PdfReader(file)
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pdf_text += text
            else:
                # Use OCR to extract text from image
                images = convert_from_path(file.name)
                for image in images:
                    pdf_text += pytesseract.image_to_string(image)
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
    return pdf_text

def process_docx(file):
    doc = Document(file)
    doc_text = ""
    for para in doc.paragraphs:
        doc_text += para.text + "\n"
    return doc_text

def process_pptx(file):
    ppt = Presentation(file)
    ppt_text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text += shape.text + "\n"
    return ppt_text

def process_and_store_files(files):
    texts = []
    metadatas = []
    for file in files:
        if file.name not in st.session_state["processed_files"]:
            if file.type == "application/pdf":
                file_text = process_pdf(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                file_text = process_docx(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                file_text = process_pptx(file)
            else:
                st.error(f"Unsupported file type: {file.type}")
                continue

            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=50)
            file_texts = text_splitter.split_text(file_text)
            if not file_texts:
                st.error(f"No text found in {file.name}")
                continue
            texts.extend(file_texts)
            file_metadatas = [{"source": f"{i}-{file.name}"} for i in range(len(file_texts))]
            metadatas.extend(file_metadatas)

            st.session_state["processed_files"][file.name] = {
                "texts": file_texts,
                "metadatas": file_metadatas
            }
        else:
            data = st.session_state["processed_files"][file.name]
            texts.extend(data["texts"])
            metadatas.extend(data["metadatas"])

    return texts, metadatas

def initialize_chain(texts, metadatas):
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    docsearch = Chroma.from_texts(texts, embeddings, metadatas=metadatas)
    message_history = ChatMessageHistory()
    memory = ConversationBufferMemory(
        memory_key="chat_history",
        output_key="answer",
        chat_memory=message_history,
        return_messages=True,
    )
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm_groq,
        chain_type="stuff",
        retriever=docsearch.as_retriever(),
        memory=memory,
        return_source_documents=True,
    )
    return chain

# Initialize session state if not already done
if "chain" not in st.session_state:
    st.session_state["chain"] = None
if "uploaded_files" not in st.session_state:
    st.session_state["uploaded_files"] = set()
if "processed_files" not in st.session_state:
    st.session_state["processed_files"] = {}
if "image_displayed" not in st.session_state:
    st.session_state["image_displayed"] = False

# Add custom CSS for better UI
st.markdown("""
    <style>
        .stTextArea textarea {
            font-size: 14px;
        }
        .stButton button {
            font-size: 14px;
        }
        .stFileUploader label {
            font-size: 14px;
        }
        .stSidebar .css-1d391kg p {
            font-size: 14px;
        }
        .small-image {
            width: 50%;
        }
    </style>
""", unsafe_allow_html=True)

st.title("PDF Conversational AI")
st.write("Upload PDF, Word, or PowerPoint files and ask questions about their content.")

uploaded_files = st.file_uploader("Choose files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

if uploaded_files:
    new_files = [file for file in uploaded_files if file.name not in st.session_state["uploaded_files"]]
    if new_files:
        with st.spinner('Processing files...'):
            texts, metadatas = process_and_store_files(new_files)
            if texts:
                if "texts" not in st.session_state:
                    st.session_state["texts"] = []
                if "metadatas" not in st.session_state:
                    st.session_state["metadatas"] = []
                st.session_state["texts"].extend(texts)
                st.session_state["metadatas"].extend(metadatas)
                st.session_state["chain"] = initialize_chain(st.session_state["texts"], st.session_state["metadatas"])
                for file in new_files:
                    st.session_state["uploaded_files"].add(file.name)
                    c.execute("INSERT INTO uploaded_files (file_name) VALUES (?)", (file.name,))
                conn.commit()
                st.success("Files processed. You can now ask questions.")
                # Display the image only once
                if not st.session_state["image_displayed"]:
                    st.image("pic.jpg", width=300, caption="Processing complete")
                    st.session_state["image_displayed"] = True
    else:
        texts = []
        metadatas = []
        for file_name, data in st.session_state["processed_files"].items():
            texts.extend(data["texts"])
            metadatas.extend(data["metadatas"])

        if "texts" not in st.session_state:
            st.session_state["texts"] = texts
        if "metadatas" not in st.session_state:
            st.session_state["metadatas"] = metadatas

        st.session_state["chain"] = initialize_chain(st.session_state["texts"], st.session_state["metadatas"])

st.write("### Ask a question about your uploaded files:")
user_input = st.text_area("Type your message here...", height=100)
if st.session_state["chain"]:
    if st.button("Send"):
        with st.spinner('Fetching the answer...'):
            chain = st.session_state["chain"]
            try:
                res = chain(user_input)
                answer = res["answer"]
                source_documents = res["source_documents"]
                if not answer:
                    general_response = llm_groq.agenerate({"prompt": user_input})
                    answer = general_response['choices'][0]['message']['content']

                # Save the question and answer to session state
                if "history" not in st.session_state:
                    st.session_state["history"] = []
                st.session_state["history"].append({"question": user_input, "answer": answer})

                # Display the question and answer
                st.write(f"**Q:** {user_input}")
                st.write(f"**A:** {answer}")

                # Save the question and answer to the database
                c.execute("INSERT INTO chat_history (question, answer) VALUES (?, ?)", (user_input, answer))
                conn.commit()
            except Exception as e:
                st.error(f"Failed to fetch the answer: {e}")

# Sidebar for Chat History
st.sidebar.write("### Chat History")
if "history" in st.session_state:
    for entry in st.session_state["history"]:
        st.sidebar.write(f"**Q:** {entry['question']}")
        st.sidebar.write(f"**A:** {entry['answer']}")
